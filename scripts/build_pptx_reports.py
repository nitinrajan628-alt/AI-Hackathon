"""Generate professional PowerPoint slide decks from the structured report data.

Reads slides from reserve_review.db (same source as build_reports.py) and
produces one .pptx per review plus pre-rendered PNG images for web display.

Run after build_reports.py:  python scripts/build_pptx_reports.py

Requires: python-pptx, comtypes (Windows, for PNG export via PowerPoint COM)
"""
from __future__ import annotations

import json
import os
import shutil
import sqlite3
import sys
import tempfile

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

DB_PATH = os.path.join(os.path.dirname(__file__), "..", "data", "reserve_review.db")
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "data", "reports")

BRAND_NAVY = RGBColor(0x1C, 0x5C, 0xAB)
BRAND_DARK = RGBColor(0x14, 0x14, 0x12)
BRAND_INK2 = RGBColor(0x52, 0x51, 0x4E)
BRAND_MUTED = RGBColor(0x89, 0x87, 0x81)
BRAND_ACCENT_SOFT = RGBColor(0xE8, 0xEF, 0xF8)
BRAND_SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
BRAND_HAIRLINE = RGBColor(0xE3, 0xE2, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CHART_PALETTE = [
    RGBColor(0x1C, 0x5C, 0xAB),
    RGBColor(0x2F, 0x8F, 0x4E),
    RGBColor(0xE0, 0x72, 0x2A),
    RGBColor(0xC7, 0x90, 0x00),
    RGBColor(0x6B, 0x4B, 0xC4),
    RGBColor(0xD1, 0x34, 0x4A),
]

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_LEFT = Inches(0.6)
MARGIN_RIGHT = Inches(0.6)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT


def _set_cell_text(cell, text: str, font_size: int = 10, bold: bool = False,
                   color: RGBColor = BRAND_DARK, align: PP_ALIGN = PP_ALIGN.LEFT):
    cell.text = str(text) if text is not None else ""
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _set_cell_fill(cell, color: RGBColor):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


class SlideRenderer:
    """Renders one structured slide dict to a PowerPoint slide."""

    def __init__(self, prs: Presentation, slide_data: dict, review_label: str):
        self.prs = prs
        self.data = slide_data
        self.review_label = review_label
        self.content = slide_data.get("content", {})
        self.top = Inches(0.4)

    def render(self):
        layout = self.prs.slide_layouts[6]  # Blank layout
        self.slide = self.prs.slides.add_slide(layout)

        # Background
        bg = self.slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = BRAND_SURFACE

        self._add_header_bar()
        self._add_section_label()
        self._add_title()
        self._add_headline()
        self._render_blocks()
        self._add_footer()

    def _add_header_bar(self):
        """Thin branded bar across the top."""
        shape = self.slide.shapes.add_shape(
            1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BRAND_NAVY
        shape.line.width = Pt(0)
        shape.line.fill.background()

    def _add_section_label(self):
        txbox = self.slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(0.25), Inches(4), Inches(0.3))
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.text = self.data.get("section", "").upper()
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = BRAND_NAVY
        p.font.name = "Segoe UI"
        self.top = Inches(0.55)

    def _add_title(self):
        txbox = self.slide.shapes.add_textbox(
            MARGIN_LEFT, self.top, CONTENT_WIDTH, Inches(0.55))
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.text = self.data.get("title", "")
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = BRAND_DARK
        p.font.name = "Segoe UI"
        self.top += Inches(0.6)

    def _add_headline(self):
        headline = self.content.get("headline")
        if not headline or headline == self.data.get("title"):
            return
        txbox = self.slide.shapes.add_textbox(
            MARGIN_LEFT, self.top, CONTENT_WIDTH, Inches(0.45))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = headline
        p.font.size = Pt(13)
        p.font.italic = True
        p.font.color.rgb = BRAND_INK2
        p.font.name = "Segoe UI"
        self.top += Inches(0.5)

    def _render_blocks(self):
        blocks = self.content.get("blocks", [])
        for block in blocks:
            if self.top > Inches(6.5):
                break
            btype = block.get("type")
            if btype == "metrics":
                self._block_metrics(block.get("items", []))
            elif btype in ("bullets", "numbered"):
                self._block_bullets(block.get("items", []), btype == "numbered")
            elif btype == "table":
                self._block_table(block.get("columns", []), block.get("rows", []))
            elif btype == "chart":
                self._block_chart(block)
            elif btype == "text":
                self._block_text(block.get("text", ""))

    def _block_metrics(self, items: list[dict]):
        if not items:
            return
        count = len(items)
        tile_width = min(Inches(3.2), CONTENT_WIDTH / count)
        tile_height = Inches(1.0)

        for i, item in enumerate(items):
            left = MARGIN_LEFT + tile_width * i + Inches(0.05) * i
            # Tile background
            shape = self.slide.shapes.add_shape(
                1, left, self.top, tile_width - Inches(0.05), tile_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = BRAND_ACCENT_SOFT
            shape.line.color.rgb = BRAND_HAIRLINE
            shape.line.width = Pt(0.5)

            # Label
            txbox = self.slide.shapes.add_textbox(
                left + Inches(0.15), self.top + Inches(0.1),
                tile_width - Inches(0.3), Inches(0.3))
            tf = txbox.text_frame
            p = tf.paragraphs[0]
            p.text = item.get("label", "")
            p.font.size = Pt(9)
            p.font.color.rgb = BRAND_MUTED
            p.font.name = "Segoe UI"

            # Value
            txbox = self.slide.shapes.add_textbox(
                left + Inches(0.15), self.top + Inches(0.4),
                tile_width - Inches(0.3), Inches(0.5))
            tf = txbox.text_frame
            p = tf.paragraphs[0]
            p.text = item.get("value", "")
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = BRAND_NAVY
            p.font.name = "Segoe UI"

        self.top += tile_height + Inches(0.3)

    def _block_bullets(self, items: list[str], numbered: bool):
        if not items:
            return
        height = Inches(min(0.35 * len(items) + 0.15, 3.0))
        txbox = self.slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(0.15), self.top, CONTENT_WIDTH - Inches(0.3), height)
        tf = txbox.text_frame
        tf.word_wrap = True
        for i, text in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            prefix = f"{i + 1}. " if numbered else "\u2022  "
            p.text = prefix + text
            p.font.size = Pt(11)
            p.font.color.rgb = BRAND_DARK
            p.font.name = "Segoe UI"
            p.space_after = Pt(6)
            p.space_before = Pt(2)
        self.top += height + Inches(0.15)

    def _block_table(self, columns: list[str], rows: list[list]):
        if not columns or not rows:
            return
        n_rows = min(len(rows) + 1, 18)
        n_cols = len(columns)
        col_width = int(min(CONTENT_WIDTH // n_cols, Inches(3.0)))
        table_width = col_width * n_cols
        row_height = Inches(0.35)
        table_height = row_height * n_rows

        shape = self.slide.shapes.add_table(
            n_rows, n_cols, MARGIN_LEFT, self.top, table_width, table_height)
        table = shape.table

        # Style header row
        for ci, col_name in enumerate(columns):
            cell = table.cell(0, ci)
            _set_cell_text(cell, col_name, font_size=9, bold=True, color=WHITE)
            _set_cell_fill(cell, BRAND_NAVY)

        # Data rows with alternating shading
        for ri, row in enumerate(rows[:n_rows - 1]):
            for ci, val in enumerate(row):
                cell = table.cell(ri + 1, ci)
                is_number = isinstance(val, (int, float))
                text = f"{val:,.1f}" if isinstance(val, float) else str(val) if val is not None else ""
                align = PP_ALIGN.RIGHT if is_number else PP_ALIGN.LEFT
                _set_cell_text(cell, text, font_size=10, color=BRAND_DARK, align=align)
                if ri % 2 == 1:
                    _set_cell_fill(cell, RGBColor(0xF5, 0xF5, 0xF2))

        self.top += table_height + Inches(0.25)

    def _block_chart(self, block: dict):
        chart_type_str = block.get("chart_type", "bar")
        x_labels = block.get("x", [])
        series_data = block.get("series", [])
        if not x_labels or not series_data:
            return

        chart_height = Inches(2.8)
        chart_width = min(CONTENT_WIDTH, Inches(10))

        chart_data = CategoryChartData()
        chart_data.categories = x_labels
        for s in series_data:
            chart_data.add_series(s.get("name", ""), s.get("values", []))

        if chart_type_str in ("grouped_bar", "waterfall"):
            ct = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif chart_type_str == "bar":
            ct = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif chart_type_str == "line":
            ct = XL_CHART_TYPE.LINE
        else:
            ct = XL_CHART_TYPE.COLUMN_CLUSTERED

        chart_frame = self.slide.shapes.add_chart(
            ct, MARGIN_LEFT, self.top, chart_width, chart_height, chart_data)
        chart = chart_frame.chart
        chart.has_legend = len(series_data) > 1

        # Apply brand colours to series
        for i, series in enumerate(chart.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = CHART_PALETTE[i % len(CHART_PALETTE)]

        # Style the value axis
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.color.rgb = BRAND_HAIRLINE
        value_axis.format.line.color.rgb = BRAND_HAIRLINE

        # Style category axis
        cat_axis = chart.category_axis
        cat_axis.format.line.color.rgb = BRAND_HAIRLINE

        self.top += chart_height + Inches(0.2)

    def _block_text(self, text: str):
        if not text:
            return
        txbox = self.slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(0.1), self.top,
            CONTENT_WIDTH - Inches(0.2), Inches(0.8))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = BRAND_INK2
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)
        self.top += Inches(0.6)

    def _add_footer(self):
        # Divider line
        shape = self.slide.shapes.add_shape(
            1, MARGIN_LEFT, Inches(7.05), CONTENT_WIDTH, Inches(0.015))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BRAND_HAIRLINE
        shape.line.width = Pt(0)
        shape.line.fill.background()

        # Footer text
        txbox = self.slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(7.1), Inches(6), Inches(0.3))
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{self.review_label} Reserve Review"
        p.font.size = Pt(8)
        p.font.color.rgb = BRAND_MUTED
        p.font.name = "Segoe UI"

        # Slide number
        txbox = self.slide.shapes.add_textbox(
            SLIDE_WIDTH - MARGIN_RIGHT - Inches(2), Inches(7.1),
            Inches(2), Inches(0.3))
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Slide {self.data.get('slide_number', '')}"
        p.font.size = Pt(8)
        p.font.color.rgb = BRAND_MUTED
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.RIGHT


def _build_title_slide(prs: Presentation, review_label: str, valuation_date: str):
    """Professional title slide with branding."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Full navy background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_NAVY

    # Company name
    txbox = slide.shapes.add_textbox(MARGIN_LEFT, Inches(1.5), Inches(10), Inches(0.5))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = "DEMO INSURANCE GROUP"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"

    # Main title
    txbox = slide.shapes.add_textbox(MARGIN_LEFT, Inches(2.5), Inches(10), Inches(1.2))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = review_label
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = "Reserve Review"
    p2.font.size = Pt(36)
    p2.font.bold = False
    p2.font.color.rgb = RGBColor(0xBB, 0xD4, 0xF0)
    p2.font.name = "Segoe UI"

    # Subtitle
    txbox = slide.shapes.add_textbox(MARGIN_LEFT, Inches(4.2), Inches(10), Inches(0.8))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Quarterly actuarial reserve review"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xBB, 0xD4, 0xF0)
    p.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = f"Valuation date: {valuation_date}"
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(0xBB, 0xD4, 0xF0)
    p2.font.name = "Segoe UI"

    # Bottom accent bar
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(7.2), SLIDE_WIDTH, Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x14, 0x3D, 0x7A)
    shape.line.width = Pt(0)
    shape.line.fill.background()


def build_pptx(review_id: str, review_label: str, valuation_date: str,
               slides: list[dict], output_path: str) -> None:
    """Build a complete .pptx deck for one review and save to output_path."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _build_title_slide(prs, review_label, valuation_date)

    for slide_data in slides:
        renderer = SlideRenderer(prs, slide_data, review_label)
        renderer.render()

    prs.save(output_path)


def render_pptx_to_png(pptx_path: str, output_dir: str) -> list[str]:
    """Render .pptx slides to PNG using PowerPoint COM automation (Windows).

    All COM interactions use a short temp path to avoid failures with long
    or space-containing paths (e.g. OneDrive sync folders). PNGs are exported
    to the temp dir then moved to the final output location.

    Returns list of PNG file paths. Falls back gracefully if PowerPoint
    is not available.
    """
    png_paths = []
    tmp_dir = None
    try:
        import comtypes.client

        abs_output = os.path.abspath(output_dir)
        os.makedirs(abs_output, exist_ok=True)

        # Use a short temp path for ALL COM file operations
        tmp_dir = tempfile.mkdtemp(dir=tempfile.gettempdir(), prefix="px")
        tmp_pptx = os.path.join(tmp_dir, "d.pptx")
        shutil.copy2(os.path.abspath(pptx_path), tmp_pptx)

        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1

        presentation = powerpoint.Presentations.Open(tmp_pptx, WithWindow=False)
        for i, slide in enumerate(presentation.Slides, start=1):
            # Export to temp dir first, then move to final location
            tmp_png = os.path.join(tmp_dir, f"s{i:02d}.png")
            slide.Export(tmp_png, "PNG", 1920, 1080)
            final_png = os.path.join(abs_output, f"slide_{i:02d}.png")
            shutil.move(tmp_png, final_png)
            png_paths.append(final_png)
        presentation.Close()
        powerpoint.Quit()
    except Exception as exc:
        print(f"  WARNING: PNG export failed ({exc}). "
              f"Install PowerPoint or run on Windows with Office.")
    finally:
        if tmp_dir and os.path.exists(tmp_dir):
            shutil.rmtree(tmp_dir, ignore_errors=True)
    return png_paths


def build_manifest(review_id: str, review_label: str, slides: list[dict],
                   png_dir: str) -> dict:
    """Build a manifest JSON for quick loading by the web app."""
    entries = []
    for slide_data in slides:
        num = slide_data["slide_number"]
        png_file = f"slide_{num + 1:02d}.png"  # +1 because title slide is slide 1
        entries.append({
            "slide_number": num,
            "section": slide_data["section"],
            "title": slide_data["title"],
            "headline": slide_data.get("content", {}).get("headline", ""),
            "png_file": png_file,
            "plain_text": slide_data.get("plain_text", ""),
        })
    return {
        "review_id": review_id,
        "review_label": review_label,
        "slide_count": len(entries) + 1,  # +1 for title slide
        "slides": entries,
    }


def main() -> int:
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    con = sqlite3.connect(DB_PATH)
    con.row_factory = sqlite3.Row
    reviews = con.execute("SELECT * FROM review_period ORDER BY sequence_no").fetchall()

    for review in reviews:
        rid = review["review_id"]
        label = review["quarter_label"]
        val_date = review["valuation_date"]

        # Load slides from DB (same data that build_reports.py wrote)
        rows = con.execute(
            "SELECT slide_number, section, title, content_json, plain_text, "
            "tags_json FROM report_slide WHERE review_id=? ORDER BY slide_number",
            (rid,)).fetchall()
        slides = []
        for r in rows:
            d = dict(r)
            d["content"] = json.loads(d.pop("content_json"))
            d["tags"] = json.loads(d.pop("tags_json"))
            slides.append(d)

        if not slides:
            print(f"  {rid}: no slides in DB, skipping")
            continue

        # Build .pptx
        review_dir = os.path.join(OUTPUT_DIR, rid)
        os.makedirs(review_dir, exist_ok=True)
        pptx_path = os.path.join(review_dir, f"{rid}.pptx")
        build_pptx(rid, label, val_date, slides, pptx_path)
        file_size = os.path.getsize(pptx_path) // 1024
        print(f"  {rid}: wrote {pptx_path} ({file_size} KB)")

        # Render to PNG
        png_paths = render_pptx_to_png(pptx_path, review_dir)
        if png_paths:
            print(f"  {rid}: exported {len(png_paths)} PNG slides")
        else:
            print(f"  {rid}: PNG export skipped (no PowerPoint COM available)")

        # Write manifest
        manifest = build_manifest(rid, label, slides, review_dir)
        manifest_path = os.path.join(review_dir, "manifest.json")
        with open(manifest_path, "w") as f:
            json.dump(manifest, f, indent=2)

    con.close()
    print("Done.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
